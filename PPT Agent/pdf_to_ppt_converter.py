from pptx import Presentation
from pptx.util import Inches
import fitz 
import os
import json
import logging
import requests
from dotenv import load_dotenv
from tool_registry import ToolRegistry
from base_tool import BaseTool
from image_fetcher import ImageFetcherTool

logging.basicConfig(level=logging.INFO)

class PDFToPPTConverterTool(BaseTool):
    def __init__(self):
        load_dotenv()  
        try:
            from config import Config
            self.api_key = os.getenv("OPENAI_API_KEY") or Config.OPENAI_API_KEY
        except Exception:
            self.api_key = os.getenv("OPENAI_API_KEY")
        if not self.api_key:
            raise Exception("OPENAI_API_KEY is not set.")

    def run(self, pdf_path):
        logging.info(f"📥 PDF received: {pdf_path}")
        logging.info(f"📄 Converting PDF to PPT: {pdf_path}")
        prs = Presentation()

        try:
            # Step 1: Extract text from the PDF using PyMuPDF
            text_content = self.extract_text(pdf_path)
            logging.info(f"Extracted text length: {len(text_content)}")

            # Step 2: Use the LLM to generate a structured summary for exactly 3 slides.
            slides_data = self.generate_slides_summary(text_content)
            logging.info(f"Generated slides data: {slides_data}")

            # Step 3: Create one slide per generated slide object.
            for slide_obj in slides_data:
                self.add_slide(prs, slide_obj['title'], slide_obj['content'], slide_obj.get('image_query', slide_obj['title']))

            # Step 4: Save the PPTX file.
            ppt_path = pdf_path.replace(".pdf", "_converted.pptx")
            prs.save(ppt_path)
            logging.info(f"✅ PPT saved successfully at: {ppt_path}")
            return ppt_path

        except Exception as e:
            logging.error(f"❌ Error during PDF to PPT conversion: {e}", exc_info=True)
            return None

    def extract_text(self, pdf_path):
        """Extracts text from all pages of the PDF using PyMuPDF."""
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text

    def generate_slides_summary(self, text):
        """
        Calls OpenAI's chat completions endpoint to generate a JSON array of exactly 3 slide objects.
        Each object must include:
          - "title": a concise slide title,
          - "content": a brief summary (max ~30 words),
          - "image_query": a term for fetching a representative image (if omitted, defaults to the title).
        """
        api_url = "https://api.openai.com/v1/chat/completions"
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        prompt = (
            "You are a PowerPoint presentation specialist. Summarize the following content into exactly 3 slides. "
            "For each slide, output a JSON object with keys 'title', 'content', and 'image_query'. "
            "The 'title' should be a concise slide title; 'content' a brief summary (max ~30 words); "
            "and 'image_query' a term suitable for fetching a representative image. "
            "Return only a valid JSON array of 3 objects with no markdown formatting or extra commentary.\n\n"
            f"Content: {text}"
        )
        data = {
            "model": "gpt-4o-mini",
            "messages": [{"role": "user", "content": prompt}],
            "max_tokens": 300,  
            "temperature": 0.3
        }
        response = requests.post(api_url, json=data, headers=headers).json()
        try:
            slides = response["choices"][0]["message"]["content"].strip()
            if slides.startswith("```"):
                slides = slides.strip("```").strip()
            # Ensure the JSON is complete
            if not slides.endswith("]"):
                slides += "]"
            slides_data = json.loads(slides)
            if isinstance(slides_data, list) and len(slides_data) == 3:
                return slides_data
            else:
                logging.error(f"Generated slides data is not a list of 3 objects: {slides_data}")
                raise Exception("Invalid slides data")
        except Exception as e:
            logging.error(f"Unexpected response from OpenAI: {response}")
            raise Exception(f"Unexpected response from OpenAI: {response}")

    def add_slide(self, prs, title, content, image_query):
        """
        Adds a slide with a title, content, and an image.
        The slide layout places the title at the top, the content on the left, and the image on the right.
        """
        # Create a blank slide layout (use layout index 6; adjust if necessary)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Add title at the top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_box.text_frame.text = title
        # Add content on the left half
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
        content_box.text_frame.text = content
        # Fetch and add an image on the right half
        image_path = self.fetch_and_save_image(image_query)
        if image_path:
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4), height=Inches(4))

    def fetch_and_save_image(self, query):
        """
        Uses the ImageFetcherTool to fetch an image based on the given query.
        """
        fetcher = ImageFetcherTool()
        return fetcher.run(query)

# Register the tool in the registry
ToolRegistry.register("pdf_to_ppt_converter", PDFToPPTConverterTool())
logging.info("✅ pdf_to_ppt_converter tool registered successfully.")

if __name__ == "__main__":
    # For standalone testing, replace with your actual PDF file path.
    converter = PDFToPPTConverterTool()
    converter.run("path/to/your/pdf.pdf")
