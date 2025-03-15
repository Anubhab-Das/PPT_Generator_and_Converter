from pptx import Presentation
from pptx.util import Inches
import os
import logging
from config import Config
from base_tool import BaseTool

logging.basicConfig(level=logging.INFO)

class PPTGeneratorTool(BaseTool):
    def run(self, main_topic, subtopics, generated_text):
        logging.info(f"✅ Generating PPT for topic: {main_topic} with subtopics: {subtopics}")
        output_path = os.path.join(Config.OUTPUT_PATH, "generated_presentation.pptx")
        logging.info(f"📁 Output Path: {output_path}")
        
        try:
            prs = Presentation()
            background_path = os.path.join(Config.ASSETS_PATH, "light_blue_gradient.png")
            
            # -------------------------------
            # Create Title Slide with Background
            # -------------------------------
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
            title_slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add the background image covering entire slide
            if os.path.exists(background_path):
                bg_shape = title_slide.shapes.add_picture(
                    background_path, 
                    0, 0, 
                    width=prs.slide_width, 
                    height=prs.slide_height
                )
                # Move background to the back by reordering the shape tree
                spTree = title_slide.shapes._spTree
                spTree.remove(bg_shape._element)
                spTree.insert(2, bg_shape._element)
            else:
                logging.warning("Background image not found; skipping background for title slide.")
            
            # Add main topic image (if available) near the top
            main_image_path = self.fetch_and_save_image(main_topic)
            if main_image_path:
                logging.info(f"🖼️ Adding main topic image for {main_topic} from {main_image_path}")
                # Place the image at (1", 0.5") with height fixed at 2.5"
                title_slide.shapes.add_picture(
                    main_image_path, 
                    Inches(1), Inches(0.5), 
                    width=Inches(5), 
                    height=Inches(2.5)
                )
                # Set title textbox below the image (image top + image height + margin)
                title_top = Inches(0.5 + 2.5 + 0.5)
            else:
                title_top = Inches(1)
            
            # Add title textbox below the image
            title_box = title_slide.shapes.add_textbox(Inches(1), title_top, Inches(8), Inches(1.5))
            title_box.text_frame.text = main_topic
            
            # -------------------------------
            # Create Subtopic Slides with Background
            # -------------------------------
            for subtopic in subtopics:
                # Retrieve generated slide data for the subtopic (expected as a list of slide dicts)
                subtopic_data_list = generated_text.get(subtopic, [])
                logging.info(f"DEBUG: For subtopic '{subtopic}', raw generated data: {subtopic_data_list}")
                
                if isinstance(subtopic_data_list, list) and subtopic_data_list:
                    for entry in subtopic_data_list:
                        # Create a blank slide for custom design
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        
                        # Add background to the slide
                        if os.path.exists(background_path):
                            bg_shape = slide.shapes.add_picture(
                                background_path, 
                                0, 0, 
                                width=prs.slide_width, 
                                height=prs.slide_height
                            )
                            spTree = slide.shapes._spTree
                            spTree.remove(bg_shape._element)
                            spTree.insert(2, bg_shape._element)
                        else:
                            logging.warning("Background image not found; skipping background for slide.")
                        
                        # Add title at top (spanning full width)
                        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
                        title_box.text_frame.text = entry.get("title", subtopic)
                        
                        # Add content on the left side below title
                        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(4))
                        content_box.text_frame.text = entry.get("content", "No content found.")
                        
                        # Add image on the right side ensuring no overlap
                        image_query = entry.get("image_query", entry.get("title", subtopic))
                        image_path = self.fetch_and_save_image(image_query)
                        if image_path:
                            logging.info(f"🖼️ Adding image for {subtopic} from {image_path}")
                            slide.shapes.add_picture(
                                image_path, 
                                Inches(5.5), Inches(1.5), 
                                width=Inches(4), 
                                height=Inches(4)
                            )
                else:
                    logging.warning(f"No slide data found for subtopic: {subtopic}")
            
            prs.save(output_path)
            logging.info(f"✅ PPT saved successfully at: {output_path}")
            return output_path
        
        except Exception as e:
            logging.error(f"❌ PPT Generation Failed: {e}")
            return None

    def fetch_and_save_image(self, query):
        from image_fetcher import ImageFetcherTool
        fetcher = ImageFetcherTool()
        return fetcher.run(query)

from tool_registry import ToolRegistry
ToolRegistry.register("ppt_generator", PPTGeneratorTool())
logging.info("✅ ppt_generator tool registered successfully.")
